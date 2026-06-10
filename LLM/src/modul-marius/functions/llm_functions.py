import json
import os
import sys
import time
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeout
from typing import Dict, Tuple

import pdfplumber
import pybreaker
from langdetect import detect as langdetect_detect, LangDetectException
from dotenv import load_dotenv
from google import genai
from google.genai import errors as genai_errors
from google.genai import types as genai_types
from pydantic import ValidationError
from supabase import create_client, Client
from tenacity import (
    retry,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
)

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
import cache as llm_cache
import supabase_logger
from schemas import (
    AnswerQuestionInput,
    AnswerQuestionOutput,
    GenerateSummaryInput,
    GenerateSummaryOutput,
)

load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", ".env"))

GEMINI_MODEL = "gemini-2.5-flash-lite"
EMBEDDING_MODEL = os.getenv("GEMINI_EMBEDDING_MODEL", "gemini-embedding-001")
EMBEDDING_DIMS = 384
TEMPERATURE = 0.7
TIMEOUT_S = 30

logger = logging.getLogger("llm-functions")

raw_api_key = os.getenv("GEMINI_API_KEY", "").strip().strip("'").strip('"')
_client = genai.Client(api_key=raw_api_key)

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_KEY")

if SUPABASE_URL and SUPABASE_KEY:
    supabase_client: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
else:
    logger.warning("Supabase URL sau KEY lipsesc. Funcționalitățile Supabase vor fi dezactivate.")
    supabase_client = None

# Se închide după 5 erori consecutive, se resetează după 60s
_breaker = pybreaker.CircuitBreaker(fail_max=5, reset_timeout=60)


# ── LLM call cu timeout + circuit breaker + retry ────────

@_breaker
def _raw_gemini(prompt: str) -> Tuple[str, Dict[str, int]]:
    resp = _client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
    usage: Dict[str, int] = {}
    if hasattr(resp, "usage_metadata") and resp.usage_metadata:
        usage = {
            "prompt_tokens": resp.usage_metadata.prompt_token_count or 0,
            "response_tokens": resp.usage_metadata.candidates_token_count or 0,
            "total_tokens": resp.usage_metadata.total_token_count or 0,
        }
    return resp.text or "", usage


@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=2, max=10),
    retry=retry_if_exception_type((TimeoutError, pybreaker.CircuitBreakerError, OSError, genai_errors.ServerError)),
)
def _call(prompt: str, function_name: str = "unknown") -> str:
    """Apel Gemini cu cache, timeout 30s, max 2 retry-uri, circuit breaker și logging Supabase."""
    key = llm_cache.make_key(prompt, GEMINI_MODEL, TEMPERATURE)
    cached = llm_cache.get(key)
    if cached is not None:
        supabase_logger.log_llm_call(
            function_name=function_name,
            model=GEMINI_MODEL,
            prompt_tokens=0,
            response_tokens=0,
            total_tokens=0,
            cached=True,
        )
        return cached

    start = time.time()
    with ThreadPoolExecutor(max_workers=1) as pool:
        future = pool.submit(_raw_gemini, prompt)
        try:
            result, usage = future.result(timeout=TIMEOUT_S)
        except FutureTimeout:
            raise TimeoutError(f"Gemini nu a răspuns în {TIMEOUT_S}s")
    duration_ms = int((time.time() - start) * 1000)

    llm_cache.set(key, result, GEMINI_MODEL)
    supabase_logger.log_llm_call(
        function_name=function_name,
        model=GEMINI_MODEL,
        prompt_tokens=usage.get("prompt_tokens", 0),
        response_tokens=usage.get("response_tokens", 0),
        total_tokens=usage.get("total_tokens", 0),
        cached=False,
        duration_ms=duration_ms,
    )
    return result


# ── PDF / vector DB helpers ───────────────────────────────

def extract_text_from_pdf(pdf_path: str) -> str:
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    raise ValueError(f"Format nesupорtat: {ext}")


def detect_language(text: str) -> str:
    """Detectează limba PDF-ului automat. Returnează codul ISO (ex: 'ro', 'en', 'fr')."""
    try:
        return langdetect_detect(text[:3000])
    except LangDetectException:
        return "ro"


def chunk_text(text: str, max_words: int = 200, overlap_words: int = 30) -> list:
    """Chunking inteligent: paragraf → propoziție → cuvinte ca fallback.
    Păstrează contextul semantic nealterat față de tăierea brutală la N cuvinte.
    """
    import re

    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]

    raw_chunks: list[str] = []

    for para in paragraphs:
        words = para.split()
        if not words:
            continue
        if len(words) <= max_words:
            raw_chunks.append(para)
        else:
            # Paragraful e prea lung — împarte în propoziții
            sentences = re.split(r"(?<=[.!?])\s+", para)
            current: list[str] = []
            for sent in sentences:
                sent_words = sent.split()
                if len(current) + len(sent_words) <= max_words:
                    current.extend(sent_words)
                else:
                    if current:
                        raw_chunks.append(" ".join(current))
                    if len(sent_words) <= max_words:
                        current = sent_words
                    else:
                        # Propoziție uriașă — fallback word split
                        for i in range(0, len(sent_words), max_words):
                            raw_chunks.append(" ".join(sent_words[i : i + max_words]))
                        current = []
            if current:
                raw_chunks.append(" ".join(current))

    if not raw_chunks:
        return raw_chunks

    # Aplică overlap semantic între chunk-uri consecutive
    if overlap_words <= 0 or len(raw_chunks) == 1:
        return raw_chunks

    overlapped = [raw_chunks[0]]
    for i in range(1, len(raw_chunks)):
        prev_words = raw_chunks[i - 1].split()
        tail = " ".join(prev_words[-overlap_words:])
        overlapped.append(tail + " " + raw_chunks[i])
    return overlapped


def store_in_vector_db(chunks: list, pdf_id: str):
    embeddings = []
    for chunk in chunks:
        resp = _client.models.embed_content(
            model=EMBEDDING_MODEL,
            contents=chunk,
            config=genai_types.EmbedContentConfig(output_dimensionality=EMBEDDING_DIMS),
        )
        embeddings.append(resp.embeddings[0].values)

    data = []
    for i, chunk in enumerate(chunks):
        data.append({
            "pdf_id": pdf_id,
            "content": chunk,
            "embedding": embeddings[i]
        })
        
    supabase_client.table("document_chunks").insert(data).execute()
    return True


def query_relevant_chunks(query: str, pdf_id: str, n_results: int = 5) -> list:
    resp = _client.models.embed_content(
        model=EMBEDDING_MODEL,
        contents=query,
        config=genai_types.EmbedContentConfig(output_dimensionality=EMBEDDING_DIMS),
    )
    query_embedding = resp.embeddings[0].values

    response = supabase_client.rpc(
        "match_document_chunks",
        {
            "query_embedding": query_embedding,
            "match_count": n_results,
            "filter_pdf_id": pdf_id
        }
    ).execute()
    
    return [item["content"] for item in response.data]


def delete_pdf_from_rag(pdf_id: str) -> None:
    """Sterge fragmentele din pgvector asociate unui PDF (garbage collection vectorial)."""
    try:
        supabase_client.table("document_chunks").delete().eq("pdf_id", pdf_id).execute()
    except Exception:
        pass


def load_pdf_into_rag(pdf_path: str, pdf_id: str) -> tuple:
    text = extract_text_from_file(pdf_path)
    language = detect_language(text)
    chunks = chunk_text(text)
    store_in_vector_db(chunks, pdf_id)
    return chunks, language


# ── Funcții publice cu contract Pydantic ──────────────────

def answer_question(question: str, pdf_id: str, language: str = "ro") -> str:
    inp = AnswerQuestionInput(question=question, pdf_id=pdf_id)

    chunks = query_relevant_chunks(inp.question, inp.pdf_id, n_results=8)
    context = "\n\n".join(chunks)
    prompt = (
        "Raspunde in aceeasi limba in care este scris materialul de mai jos.\n"
        "Esti Asistentul Virtual InsideUGAL. Rolul tau este sa ajuti studentii sa inteleaga regulamentele si documentele administrative ale UGAL.\n"
        "Reguli:\n"
        "- Raspunde EXCLUSIV pe baza informatiilor din materialul administrativ furnizat.\n"
        "- Daca nu cunosti raspunsul pe baza documentelor, directioneaza studentul catre secretariat si fii mereu politicos.\n"
        "- Nu inventa informatii care nu sunt in material.\n"
        "- Daca materialul nu contine absolut nimic relevant pentru intrebare, spune politicos ca nu ai informatia in baza de date.\n\n"
        f"Material:\n{context}\n\nIntrebare: {inp.question}"
    )

    for attempt in range(2):
        raw = _call(prompt, function_name="answer_question")
        try:
            return AnswerQuestionOutput(answer=raw).answer
        except ValidationError as exc:
            if attempt == 1:
                raise ValueError(f"Output LLM invalid după retry: {exc}") from exc

    raise RuntimeError("unreachable")


def generate_summary(pdf_id: str, language: str = "ro") -> str:
    inp = GenerateSummaryInput(pdf_id=pdf_id)

    chunks = query_relevant_chunks(
        "rezumat general curs introducere concepte principale", inp.pdf_id, n_results=8
    )
    context = "\n\n".join(chunks)
    prompt = (
        "Esti Asistentul Virtual InsideUGAL.\n"
        "Raspunde in aceeasi limba in care este scris materialul de mai jos.\n"
        "Creeaza un rezumat structurat al acestui document administrativ/regulament, respectand EXACT acest format:\n\n"
        "## Reguli stricte si proceduri\n"
        "- [maxim 5 reguli cheie extrase din text]\n\n"
        "## Informatii utile\n"
        "- **Informatie**: explicatie scurta si clara\n"
        "- [repeta pentru fiecare detaliu important]\n\n"
        "## Pasi de urmat (daca se aplica)\n"
        "[2-3 fraze cu pasii pe care studentul trebuie sa ii urmeze conform procedurii]\n\n"
        "Reguli:\n"
        "- Foloseste limbaj simplu, fara jargon inutil\n"
        "- Fii concis — studentul trebuie sa inteleaga in 2 minute\n"
        "- Nu copia fraze din material, reformuleaza cu cuvinte proprii\n\n"
        f"Material:\n{context}"
    )

    for attempt in range(2):
        raw = _call(prompt, function_name="generate_summary")
        try:
            return GenerateSummaryOutput(summary=raw).summary
        except ValidationError as exc:
            if attempt == 1:
                raise ValueError(f"Rezumat invalid după retry: {exc}") from exc

    raise RuntimeError("unreachable")
